#!/usr/bin/env python3
# Copyright 2026 Eluu Labs — Apache-2.0 (see ../LICENSE)
"""
template_follow — inherit a user's corporate .pptx and fill it (template mode).

The point of template-following is to look NATIVE to the supplied deck. So we
(1) read the deck's theme tokens (fonts + palette) to drive the design-system
lock, (2) inventory its slides as raw material for the audit / frame-map, and
(3) clone a chosen source slide and swap its text run-by-run, preserving every
run's font, size, colour, and weight — a number changes without the styling.

Functions:
    extract_theme(path)        -> {"fonts": {...}, "colors": {...}}
    inventory(path)            -> [{"index", "layout", "shapes":[...]}, ...]
    load(path)                 -> pptx.Presentation
    clone_slide(prs, index)    -> a new slide duplicated from source `index`
    replace_text(slide, map)   -> swap exact substrings, keeping run formatting

Run `python template_follow.py corporate.pptx` to print the theme + inventory.

Note on images: clone_slide copies shapes (incl. their formatting) onto a new
slide that inherits the same layout/master, so text and theme carry over. Picture
shapes reference relationships on the source slide part; for image-heavy
templates, re-place images during fill rather than relying on the clone, or copy
the whole slide part. Text/table/shape templates clone cleanly.
"""
import sys, copy, json
from pptx import Presentation
from lxml import etree

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
RT_THEME = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"

def load(path):
    return Presentation(path)

def _theme_root(prs):
    part = prs.slide_masters[0].part.part_related_by(RT_THEME)
    el = getattr(part, "_element", None)
    return el if el is not None else etree.fromstring(part.blob)

def extract_theme(path):
    """Read the deck's font scheme + colour scheme — the inherited design system."""
    prs = Presentation(path)
    root = _theme_root(prs)
    elems = root.find(f"{A}themeElements")
    fonts = {}
    fs = elems.find(f"{A}fontScheme") if elems is not None else None
    if fs is not None:
        for key, tag in (("major", "majorFont"), ("minor", "minorFont")):
            node = fs.find(f"{A}{tag}")
            latin = node.find(f"{A}latin") if node is not None else None
            fonts[key] = latin.get("typeface") if latin is not None else None
    colors = {}
    cs = elems.find(f"{A}clrScheme") if elems is not None else None
    if cs is not None:
        for c in cs:
            name = etree.QName(c).localname
            srgb = c.find(f"{A}srgbClr"); sysc = c.find(f"{A}sysClr")
            if srgb is not None: colors[name] = "#" + srgb.get("val")
            elif sysc is not None: colors[name] = "#" + (sysc.get("lastClr") or "000000")
    return {"fonts": fonts, "colors": colors}

def inventory(path):
    """List every slide's shapes/placeholders/text — raw material for the audit."""
    prs = Presentation(path); out = []
    for i, sl in enumerate(prs.slides):
        shapes = []
        for sh in sl.shapes:
            shapes.append({
                "name": sh.name,
                "type": str(sh.shape_type),
                "is_placeholder": bool(sh.is_placeholder),
                "text": (sh.text_frame.text if sh.has_text_frame else None),
            })
        out.append({"index": i, "layout": sl.slide_layout.name, "shapes": shapes})
    return out

def clone_slide(prs, index):
    """Duplicate source slide `index` into a new slide that inherits its layout."""
    source = prs.slides[index]
    new = prs.slides.add_slide(source.slide_layout)
    # drop the placeholder shapes the layout auto-added, then copy the source tree
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in source.shapes:
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    return new

def replace_text(slide, mapping):
    """Swap exact substrings across all text runs, preserving each run's style."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                t = run.text
                for k, v in mapping.items():
                    if k in t:
                        t = t.replace(k, v)
                if t != run.text:
                    run.text = t
    return slide

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("usage: template_follow.py <deck.pptx>"); sys.exit(1)
    path = sys.argv[1]
    print("THEME:", json.dumps(extract_theme(path), indent=2))
    inv = inventory(path)
    print(f"SLIDES: {len(inv)}")
    for s in inv:
        labels = [ (sh["text"] or sh["name"])[:32] for sh in s["shapes"] ]
        print(f"  [{s['index']}] {s['layout']}: {labels}")
