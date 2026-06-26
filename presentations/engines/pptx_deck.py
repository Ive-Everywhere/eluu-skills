#!/usr/bin/env python3
# Copyright 2026 Eluu Labs — Apache-2.0 (see ../LICENSE)
"""
pptx_deck — a small, opinionated python-pptx engine for high-craft decks.

What it gives you:
  - a 12-column + 8px baseline grid on a 1280x720 px canvas (px -> EMU)
  - measured text: every string is sized against the real .ttf (Pillow) and
    pre-wrapped, so it can't overflow under font fallback
  - clean primitives: filled/outlined rects, connectors with real arrowheads,
    kicker rows, action titles, footers
  - a per-element geometry recorder that emits *.layout.json for layout_check.py
  - font EMBEDDING into the .pptx (so PowerPoint/Keynote render your fonts)

Render to PNG/PDF with LibreOffice:
    soffice --headless --convert-to pdf --outdir out deck.pptx
    pdftoppm -png -r 110 out/deck.pdf out/slide

Stage your .ttf files in a fonts/ dir next to your build script and point the
STYLE map at them. Run `python pptx_deck.py demo.pptx --layout out/layout` to
produce a 3-slide demo plus its layout JSON.
"""
import os, sys, json, zipfile
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from PIL import ImageFont
from lxml import etree

# ---- configure your fonts here -------------------------------------------
# Each style key -> (OOXML typeface name, bold flag, measurement .ttf filename).
# Use real font names; the .ttf is only used to MEASURE so text can't overflow.
FONT_DIR = os.environ.get("DECK_FONT_DIR",
                          os.path.join(os.path.dirname(os.path.abspath(__file__)), "fonts"))
STYLE = {
    "display":  ("Inter", True,  "Inter-Black.ttf"),    # cover claims
    "title":    ("Inter", True,  "Inter-Bold.ttf"),     # action titles
    "subhead":  ("Inter", False, "Inter-Medium.ttf"),
    "body":     ("Inter", False, "Inter-Regular.ttf"),
    "strong":   ("Inter", True,  "Inter-Bold.ttf"),
    "mono":     ("JetBrains Mono", False, "JetBrainsMono-Medium.ttf"),
}
# typefaces to embed into the .pptx (regular + bold slots per family)
EMBED = [("Inter", [("regular", "Inter-Regular.ttf"), ("bold", "Inter-Bold.ttf")]),
         ("JetBrains Mono", [("regular", "JetBrainsMono-Medium.ttf")])]

DISPLAY="display"; TITLE="title"; SUBHEAD="subhead"; BODY="body"; STRONG="strong"; MONO="mono"

# ---- palette (override as needed) ----------------------------------------
INK=RGBColor(0x12,0x14,0x1B); INK70=RGBColor(0x3A,0x3F,0x4C); INK40=RGBColor(0x8A,0x90,0x99)
PAPER=RGBColor(0xF5,0xF2,0xEC); HAIR=RGBColor(0xE3,0xDE,0xD3); ACCENT=RGBColor(0xE5,0x48,0x1F)
PAPER70=RGBColor(0xB9,0xB4,0xA8); INKLINE=RGBColor(0x2A,0x2D,0x36)

# ---- grid / units --------------------------------------------------------
W,H=1280,720; ML,MR,MT=72,72,64; CONTENT_L=ML; CONTENT_R=W-MR; CONTENT_W=CONTENT_R-CONTENT_L
NCOL=12; GUT=16; COLW=(CONTENT_W-(NCOL-1)*GUT)/NCOL
RULE_Y=152; CONTENT_TOP=188; FOOT_RULE=672
def colx(i): return CONTENT_L+i*(COLW+GUT)
def spanw(n): return n*COLW+(n-1)*GUT
EMU_PX=9525
def X(px): return Emu(int(round(px*EMU_PX)))
def FS(px): return Pt(px*0.75)

# ---- text measurement ----------------------------------------------------
_fc={}
def _font(key,size):
    ttf=STYLE[key][2]; k=(ttf,int(size))
    if k not in _fc: _fc[k]=ImageFont.truetype(os.path.join(FONT_DIR,ttf),int(size))
    return _fc[k]
def meas(text,key,size): return _font(key,size).getlength(text)
def wrap(text,key,size,maxw):
    out=[]; cur=""
    for w_ in text.split():
        t=(cur+" "+w_).strip()
        if meas(t,key,size)<=maxw or not cur: cur=t
        else: out.append(cur); cur=w_
    if cur: out.append(cur)
    return out or [""]
def fit_one(text,key,hi,lo,maxw):
    s=hi
    while s>lo and meas(text,key,s)>maxw: s-=1
    return s

# ---- presentation + layout recorder --------------------------------------
prs=Presentation(); prs.slide_width=X(W); prs.slide_height=X(H); BLANK=prs.slide_layouts[6]
_LAYOUTS=[]; _CUR=None; _ORD=[0]
def _hex(rgb):
    try: return "#%02x%02x%02x"%(rgb[0],rgb[1],rgb[2])
    except Exception: return None
def _new_layout():
    global _CUR; _CUR=[]; _LAYOUTS.append(_CUR); _ORD[0]=0
def _record(kind,x,y,w,h,el=None,text=None,fontSize=None,lineCount=None,fill=None):
    if _CUR is None: return
    _ORD[0]+=1
    e={"order":_ORD[0],"kind":kind,"bbox":[round(x,1),round(y,1),round(w,1),round(h,1)]}
    if el: e["name"]=el
    if text is not None: e["text"]=str(text)[:100]
    if fontSize: e["resolvedFontSize"]=round(fontSize,1)
    if lineCount: e["textLayout"]={"lineCount":lineCount}
    if fill: e["fillColor"]=fill
    _CUR.append(e)

# ---- primitives ----------------------------------------------------------
def slide(bg=PAPER):
    _new_layout()
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,X(0),X(0),X(W),X(H))
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s
def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,radius=None,shape=MSO_SHAPE.RECTANGLE,el=None):
    if radius is not None: shape=MSO_SHAPE.ROUNDED_RECTANGLE
    sp=s.shapes.add_shape(shape,X(x),X(y),X(w),X(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw*0.75)
    if radius is not None:
        try: sp.adjustments[0]=radius
        except Exception: pass
    _record("shape",x,y,w,h,el=el,fill=(_hex(fill) if fill is not None else None))
    return sp
def line(s,x1,y1,x2,y2,color,w=1.0,arrow=False,harrow=False):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,X(x1),X(y1),X(x2),X(y2))
    c.line.color.rgb=color; c.line.width=Pt(w*0.75); c.shadow.inherit=False
    ln=c.line._get_or_add_ln()
    if harrow: ln.append(ln.makeelement(qn('a:headEnd'),{'type':'triangle','w':'med','len':'med'}))
    if arrow:  ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))
    return c
def hline(s,x1,x2,y,color,w=1.0,**k): return line(s,x1,y,x2,y,color,w,**k)

def _tf(s,x,y,w,h,anchor=MSO_ANCHOR.TOP,wrap_=False):
    t=s.shapes.add_textbox(X(x),X(y),X(w),X(h)); tf=t.text_frame
    tf.word_wrap=wrap_; tf.vertical_anchor=anchor; tf.auto_size=MSO_AUTO_SIZE.NONE
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    return tf
def _r(p,txt,key,size,color,track=None,caps=False):
    name,bold,_=STYLE[key]
    r=p.add_run(); r.text=txt.upper() if caps else txt
    r.font.name=name; r.font.bold=bold; r.font.size=FS(size); r.font.color.rgb=color
    if track is not None: r.font._rPr.set('spc',str(int(track*100)))
    return r
def _alignx(x,w,tw,align):
    if align==PP_ALIGN.CENTER: return x+(w-tw)/2
    if align==PP_ALIGN.RIGHT:  return x+w-tw
    return x

def block(s,x,y,w,text,key,size,color,leading=1.3,align=PP_ALIGN.LEFT,
          accents=(),acol=ACCENT,track=None,caps=False,safety=1.0,el=None):
    """Measured, pre-wrapped text block. `accents` = phrases to colour with acol."""
    def split_runs(linetext):
        for ph in accents:
            i=linetext.find(ph)
            if i>=0:
                out=[]
                if linetext[:i]: out.append((linetext[:i],False))
                out.append((ph,True))
                if linetext[i+len(ph):]: out.append((linetext[i+len(ph):],False))
                return out
        return [(linetext,False)]
    lines=wrap(text,key,size,w*safety)
    tf=_tf(s,x,y,w+6,size*leading*len(lines)+size*0.5)
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=leading
        for t,acc in split_runs(ln): _r(p,t,key,size,(acol if acc else color),track,caps)
    tw=min(w,max((meas(ln,key,size) for ln in lines),default=0))
    _record("text",_alignx(x,w,tw,align),y,tw,size*leading*len(lines),
            el=el,text=text,fontSize=size,lineCount=len(lines))
    return y+size*leading*len(lines)

def label(s,x,y,w,text,key,size,color,align=PP_ALIGN.LEFT,track=None,caps=False,
          h=None,anchor=MSO_ANCHOR.MIDDLE,el=None):
    h=h or size*1.5
    tf=_tf(s,x,y,w,h,anchor=anchor)
    p=tf.paragraphs[0]; p.alignment=align; _r(p,text,key,size,color,track,caps)
    tw=min(w,meas(text,key,size)); gh=size*1.25
    ry=y+max(0,(h-gh)/2) if anchor==MSO_ANCHOR.MIDDLE else y
    _record("text",_alignx(x,w,tw,align),ry,tw,gh,el=el,text=text,fontSize=size,lineCount=1)
    return y+h

def kicker(s,text,dark=False,x=None,y=52):
    x=CONTENT_L if x is None else x
    rect(s,x,y+4,8,8,fill=ACCENT,el="kicker-marker")
    label(s,x+18,y,360,text,MONO,11,(PAPER70 if dark else INK70),track=2.4,caps=True,h=16,el="kicker-label")
def header(s,kick,title,accents=(),dark=False):
    base=PAPER if dark else INK
    kicker(s,kick,dark=dark)
    maxw=spanw(11); size=fit_one(title,TITLE,30,24,maxw)
    block(s,CONTENT_L,92,maxw,title,TITLE,size,base,leading=1.14,accents=accents,
          acol=ACCENT,el="action-title")
    hline(s,CONTENT_L,CONTENT_R,RULE_Y,(INKLINE if dark else HAIR),1.2)
    return CONTENT_TOP
def footer(s,num,total,dark=False,brand="ELUU"):
    col=PAPER70 if dark else INK40
    hline(s,CONTENT_L,CONTENT_R,FOOT_RULE,(INKLINE if dark else HAIR),1.0)
    label(s,CONTENT_L,680,420,brand,MONO,9,col,track=2.6,caps=True,h=16)
    label(s,CONTENT_R-200,680,200,f"{num:02d} / {total:02d}",MONO,9,col,align=PP_ALIGN.RIGHT,h=16)

# ---- font embedding (PowerPoint/Keynote honour this) ---------------------
def embed_fonts(path):
    P='http://schemas.openxmlformats.org/presentationml/2006/main'
    R='http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    CT='http://schemas.openxmlformats.org/package/2006/content-types'
    RL='http://schemas.openxmlformats.org/package/2006/relationships'
    z=zipfile.ZipFile(path,'r'); items={n:z.read(n) for n in z.namelist()}; z.close()
    ct=etree.fromstring(items['[Content_Types].xml'])
    if not any(d.get('Extension')=='fntdata' for d in ct if d.tag==f'{{{CT}}}Default'):
        d=etree.SubElement(ct,f'{{{CT}}}Default'); d.set('Extension','fntdata'); d.set('ContentType','application/x-fontdata')
    rels=etree.fromstring(items['ppt/_rels/presentation.xml.rels'])
    pres=etree.fromstring(items['ppt/presentation.xml'])
    efl=etree.Element(f'{{{P}}}embeddedFontLst'); n=0
    for tf,slots in EMBED:
        ef=etree.SubElement(efl,f'{{{P}}}embeddedFont')
        etree.SubElement(ef,f'{{{P}}}font').set('typeface',tf)
        for slot,fn in slots:
            n+=1; items[f'ppt/fonts/font{n}.fntdata']=open(os.path.join(FONT_DIR,fn),'rb').read()
            rid=f'rIdF{n}'
            rel=etree.SubElement(rels,f'{{{RL}}}Relationship')
            rel.set('Id',rid); rel.set('Type',f'{R}/font'); rel.set('Target',f'fonts/font{n}.fntdata')
            etree.SubElement(ef,f'{{{P}}}{slot}').set(f'{{{R}}}id',rid)
    kids=list(pres); pos=len(kids)
    for tag in ('notesSz','sldSz','sldIdLst'):
        idx=[j for j,c in enumerate(kids) if c.tag==f'{{{P}}}{tag}']
        if idx: pos=idx[-1]+1; break
    pres.insert(pos,efl); pres.set('embedTrueTypeFonts','1'); pres.set('saveSubsetFonts','0')
    items['[Content_Types].xml']=etree.tostring(ct,xml_declaration=True,encoding='UTF-8',standalone=True)
    items['ppt/_rels/presentation.xml.rels']=etree.tostring(rels,xml_declaration=True,encoding='UTF-8',standalone=True)
    items['ppt/presentation.xml']=etree.tostring(pres,xml_declaration=True,encoding='UTF-8',standalone=True)
    z=zipfile.ZipFile(path,'w',zipfile.ZIP_DEFLATED)
    for nm,d in items.items(): z.writestr(nm,d)
    z.close()

def save(path, layout_dir=None):
    if layout_dir:
        os.makedirs(layout_dir,exist_ok=True)
        for i,els in enumerate(_LAYOUTS):
            json.dump({"slide":{"frame":{"left":0,"top":0,"width":W,"height":H}},"elements":els},
                      open(os.path.join(layout_dir,f"slide-{i+1:02d}.layout.json"),"w"),indent=1)
    prs.save(path)
    try: embed_fonts(path)
    except Exception as e: print("embed_fonts skipped:", e)

# ---- minimal demo (3 slides) ---------------------------------------------
def demo():
    # 1) cover
    s=slide(INK)
    rect(s,CONTENT_L,150,46,5,ACCENT)
    block(s,CONTENT_L,170,spanw(11),"A deck that survives the contact-sheet test.",
          DISPLAY,46,PAPER,leading=1.05,accents=["contact-sheet test."])
    block(s,CONTENT_L,300,spanw(8),
          "Claim-led titles, one proof object per slide, an authored visual system.",
          BODY,18,PAPER70,leading=1.35)
    # 2) content slide with a tiny authored bar chart (shapes)
    s=slide(PAPER); header(s,"EXAMPLE","Authored shapes beat default chart junk.",accents=["beat default chart junk."])
    vals=[3,5,8,13]; labels=["Q1","Q2","Q3","Q4"]; bx=CONTENT_L; by=560; bw=120; gap=40; maxv=max(vals)
    for i,(v,l) in enumerate(zip(vals,labels)):
        hpx=v/maxv*300; x=bx+i*(bw+gap)
        rect(s,x,by-hpx,bw,hpx,fill=(INK if i==len(vals)-1 else ACCENT))
        label(s,x,by-hpx-26,bw,str(v),STRONG,16,INK,align=PP_ALIGN.CENTER)
        label(s,x,by+6,bw,l,BODY,13,INK70,align=PP_ALIGN.CENTER)
    footer(s,2,3)
    # 3) closing
    s=slide(INK)
    block(s,140,300,W-280,"Build it, gate it, ship it.",DISPLAY,40,PAPER,
          leading=1.05,align=PP_ALIGN.CENTER,accents=["ship it."])
    return

if __name__=="__main__":
    out=next((a for a in sys.argv[1:] if not a.startswith("--")),"demo.pptx")
    layout=None
    if "--layout" in sys.argv: layout=sys.argv[sys.argv.index("--layout")+1]
    demo(); save(out,layout)
    print("saved",out,"slides:",len(prs.slides._sldIdLst))
